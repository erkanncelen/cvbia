from reportlab.lib.pagesizes import A4
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.lib.utils import ImageReader
from PIL import Image, ImageDraw
import fitz
import io
import itertools
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm
from pypdf import PdfMerger
import sys
import textwrap
from tqdm import trange
import yaml
import re


def generate_download_file_name(yaml_input):
    cv_data = load_yaml(yaml_input)
    yaml_checker(cv_data)
    first_name = cv_data["first_name"]
    last_name = cv_data["last_name"]
    file_name_formatted = re.sub(" ", "-", f"{first_name}-{last_name}_CV")
    print(file_name_formatted)
    return file_name_formatted


def load_yaml(yaml_input):
    if yaml_input:
        return yaml.safe_load(yaml_input)
    else:
        with open("cv_data.yaml", "r") as file:
            return yaml.safe_load(file)


def cleanup_files(directories: list = ["cv_pages", "cv_images"]):
    for directory in directories:
        for filename in os.listdir(directory):
            f = os.path.join(directory, filename)
            # checking if it is a file
            if os.path.isfile(f) and filename != ".gitkeep":
                os.remove(f)


def write(
    c: object,
    x: int,
    y: int,
    text: str,
    width: int = 60,
    font: str = "regular",
    punto: int = 12,
    color: str = "dark_grey",
    spacing: int = 12,
    url: str = None,
):
    """
    Writes given text input with specified arguments.

    Arguments:
        c: canvas object
        x: x coordinate on the page
        y: y coordinate on the page
        text: text input to be written, must be a string
        font: regular, bold
        punto: font size
        color: white, trans_white, purple, dark_grey, light_grey
        spacing: line spacing, only used for mutliple line inputs
        url: URL to be written, starts with the URL domain, i.e. no "https://www."
    """
    # SET FONT
    c.setFont(font, punto)

    # SET COLOR
    if color == "white":
        c.setFillColorRGB(0.95, 0.95, 0.95, 1)
    elif color == "trans_white":
        c.setFillColorRGB(0.95, 0.95, 0.95, 0.8)
    elif color == "purple":
        c.setFillColorRGB(108 / 255, 29 / 255, 96 / 255, 1)
    elif color == "dark_grey":
        c.setFillColorRGB(0.5, 0.5, 0.5, 1)
    elif color == "light_grey":
        c.setFillColorRGB(0.3, 0.3, 0.3, 1)
    else:
        raise Exception(
            f"{color} is not available. Try: white, trans_white, purple, dark_grey, light_grey"
        )

    # DRAW STRING: for loop is used to manage multiple line inputs
    if text:
        for line in list(
            itertools.chain.from_iterable(
                [textwrap.wrap(x, width=width) for x in text.splitlines()]
            )
        ):
            c.drawString(x, y, f"{line}")
            if url is not None:
                c.linkURL(
                    "https://www." + url,
                    rect=(
                        x,
                        y,
                        x + (0.66 * stringWidth(text, font, punto + 4)),
                        y + punto,
                    ),
                    relative=0,
                    thickness=0,
                )
            y -= spacing

    return y


def draw_background(c, image=None, path="backgrounds/background1.jpg"):
    """
    Fills the background of document.
    """
    if image:
        bg_image = image
        c.drawImage(bg_image, 0, 0, width=A4[1], height=A4[0])
    elif path:
        bg_image = path
        c.drawImage(bg_image, 0, 0, width=A4[1], height=A4[0])

    return None


def draw_picture(c, image=None, path: str = "images/tom_cruise.jpg"):
    """
    Crops and draws the circular photo of the person on CV.
    """
    if image:
        im = Image.open(io.BytesIO(image))
        bigsize = (im.size[0] * 3, im.size[1] * 3)
        mask = Image.new("L", bigsize, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0) + bigsize, fill=255)
        mask = mask.resize(im.size, Image.LANCZOS)
        im.putalpha(mask)
        img_byte_arr = io.BytesIO()
        im.save(img_byte_arr, format="PNG")
        img_byte_arr = img_byte_arr.getvalue()
        pp_image = ImageReader(io.BytesIO(img_byte_arr))
        c.drawImage(pp_image, 20, 455, width=100, height=100, mask="auto")
    elif path:
        im = Image.open(path)
        bigsize = (im.size[0] * 3, im.size[1] * 3)
        mask = Image.new("L", bigsize, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0) + bigsize, fill=255)
        mask = mask.resize(im.size, Image.LANCZOS)
        im.putalpha(mask)
        img_byte_arr = io.BytesIO()
        im.save(img_byte_arr, format="PNG")
        img_byte_arr = img_byte_arr.getvalue()
        pp_image = ImageReader(io.BytesIO(img_byte_arr))
        c.drawImage(pp_image, 20, 455, width=100, height=100, mask="auto")

    return None


def page_end_checker(y: int, exp, spacing: int = 12, punto: int = 10):
    """
    This function is used to assess if the current cv page will be enough
    to display the next experience block. If not, a new page should be created.
    """

    def size_checker(y: int, text):
        if text:
            for line in text.splitlines():
                if len(text.splitlines()) > 1:
                    y -= spacing
        return y

    y = size_checker(y, text=f"{exp['title']} @ {exp['company']}")
    y -= 10
    y = size_checker(y, text=f"{exp['start']} - {exp['end']}")
    y -= 20
    y = size_checker(y, text=exp["description"])
    y -= 3
    y = size_checker(y, text=exp["technologies"])
    y -= 25
    print(y)
    if y < 0:
        return True
    else:
        return False


def merge_pdfs(input_directory: str = "cv_pages", output_directory: str = "cv_output"):
    """
    This function merges all pdfs into one, in the given directory.
    Also creates a .pptx copy.
    """
    try:
        os.remove(f"{output_directory}/cv.pdf")
        os.remove(f"{output_directory}/cv.pptx")
    except:
        pass
    pdfs = []
    for filename in os.listdir(input_directory):
        f = os.path.join(input_directory, filename)
        # checking if it is a file
        if os.path.isfile(f) and filename != ".gitkeep" and filename != "cv.pptx":
            pdfs.append(f)

    pdfs = sorted(pdfs, reverse=False)
    print(pdfs)

    merger = PdfMerger()

    for pdf in pdfs:
        merger.append(pdf)

    merger.write(f"{output_directory}/cv.pdf")
    merger.close()
    generate_pptx_from_pdf(
        pdf_file=f"{output_directory}/cv.pdf", output_file=f"{output_directory}/cv.pptx"
    )

    return None


def generate_pptx_from_pdf(
    pdf_file: str,
    output_file: str = "PDF_FILE.pptx",
    resolution: int = 300,
    start_page: int = 0,
    page_count: int = None,
):
    """
    Convert a PDF slideshow to Powerpoint PPTX.

    Source: https://github.com/kevinmcguinness/pdf2pptx/blob/master/pdf2pptx/__init__.py

    Renders each page as a PNG image and creates the resulting Powerpoint
    slideshow from these images. Useful when you want to use Powerpoint
    to present a set of PDF slides (e.g. slides from Beamer). You can then
    use the presentation capabilities of Powerpoint (notes, ink on slides,
    etc.) with slides created in LaTeX.

    Arguments:
        pdf_file: location of pdf file to convert to pptx file
        output_file: location to save the pptx (default: PDF_FILE.pptx)
        resolution: resolution in dots per inch (default: 300)
        start_page: first page in the pdf to copy to the pptx
        page_count: number of pages in the pdf to copy to the pptx
    """
    doc = fitz.open(pdf_file)

    try:
        print(pdf_file, "contains", doc.pageCount, "slides")
    except:
        print(pdf_file, "contains", doc.page_count, "slides")

    if page_count is None:
        page_count = doc.page_count

    # transformation matrix: slide to pixmap
    zoom = resolution / 72
    matrix = fitz.Matrix(zoom, zoom, 0)

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = doc.load_page(0)

    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)

    # iterate over slides
    for page_no in trange(start_page, start_page + page_count):
        page = doc.load_page(page_no)

        # write slide as a pixmap
        pixmap = page.get_pixmap(matrix=matrix)
        image_data = pixmap.tobytes("png")

        image_file = io.BytesIO(image_data)

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

    if output_file is None:
        output_file = Path(pdf_file).with_suffix(".pptx")

    # save presentation
    try:
        prs.save(output_file)
    except PermissionError as err:
        print(err, file=sys.stderr)
        sys.exit(1)


def yaml_checker(yaml):
    """
    This function is used to check whether the yaml contains the required fields.
    """

    assert yaml[
        "first_name"
    ], "'first_name' field is missing in yaml. this is a mandatory field."
    assert yaml[
        "last_name"
    ], "'last_name' field is missing in yaml. this is a mandatory field."
    assert "role" in yaml, "'role' field is missing in yaml. this is a mandatory field."
    assert (
        "email" in yaml
    ), "'email' field is missing in yaml. this is a mandatory field."

    assert (
        "about_me" in yaml
    ), "'about_me' field is missing in yaml. this is a mandatory field."
    assert (
        "education" in yaml
    ), "'education' field is missing in yaml. this is a mandatory field."
    assert (
        "biography" in yaml
    ), "'biography' field is missing in yaml. this is a mandatory field."
    assert (
        "roles" in yaml
    ), "'roles' field is missing in yaml. this is a mandatory field."
    assert (
        "certifications" in yaml
    ), "'certifications' field is missing in yaml. this is a mandatory field."
    assert (
        "competences" in yaml
    ), "'competences' field is missing in yaml. this is a mandatory field."
    assert (
        "experience" in yaml
    ), "'experience' field is missing in yaml. this is a mandatory field."

    for element in yaml["education"]:
        assert (
            "degree" in element
        ), "'degree' field is missing in one of 'education' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "institution" in element
        ), "'institution' field is missing in one of 'education' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "year" in element
        ), "'year' field is missing in one of 'education' fields. you can leave them empty but you shouldn't delete the fields."
    for element in yaml["roles"]:
        assert (
            "title" in element
        ), "'title' field is missing in one of 'roles' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "description" in element
        ), "'description' field is missing in one of 'roles' fields. you can leave them empty but you shouldn't delete the fields."
    if yaml["certifications"]:
        for element in yaml["certifications"]:
            assert (
                "title" in element
            ), "'title' field is missing in one of 'certifications' fields. you can leave them empty but you shouldn't delete the fields."
    for element in yaml["competences"]:
        assert (
            "title" in element
        ), "'title' field is missing in one of 'competences' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "description" in element
        ), "'description' field is missing in one of 'competences' fields. you can leave them empty but you shouldn't delete the fields."
    for element in yaml["experience"]:
        assert (
            "title" in element
        ), "'title' field is missing in one of 'experience' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "company" in element
        ), "'company' field is missing in one of 'experience' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "start" in element
        ), "'start' field is missing in one of 'experience' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "end" in element
        ), "'end' field is missing in one of 'experience' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "description" in element
        ), "'description' field is missing in one of 'experience' fields. you can leave them empty but you shouldn't delete the fields."
        assert (
            "technologies" in element
        ), "'technologies' field is missing in one of 'experience' fields. you can leave them empty but you shouldn't delete the fields."

    return None
